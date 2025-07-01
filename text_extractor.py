import os
import pdfplumber
import docx
from pptx import Presentation
from io import BytesIO

def extract_text_from_pdf(file_content: BytesIO) -> str:
    """Extracts text from a PDF file."""
    text = ""
    with pdfplumber.open(file_content) as pdf_reader:
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_content: BytesIO) -> str:
    """Extracts text from a DOCX file."""
    doc = docx.Document(file_content)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_content: BytesIO) -> str:
    """Extracts text from a PPTX file."""
    prs = Presentation(file_content)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_txt(file_content: BytesIO) -> str:
    """Extracts text from a TXT file."""
    return file_content.read().decode('utf-8')

def get_text_from_files(uploaded_files: list) -> str:
    """
    Extracts text from a list of uploaded files of various types.

    Args:
        uploaded_files (list): A list of Streamlit UploadedFile objects.

    Returns:
        str: The concatenated text content of all files.
    """
    full_text = ""
    for file in uploaded_files:
        file_extension = os.path.splitext(file.name)[1].lower()
        file_content = BytesIO(file.getvalue())

        if file_extension == ".pdf":
            full_text += extract_text_from_pdf(file_content)
        elif file_extension == ".docx":
            full_text += extract_text_from_docx(file_content)
        elif file_extension == ".pptx":
            full_text += extract_text_from_pptx(file_content)
        elif file_extension == ".txt":
            full_text += extract_text_from_txt(file_content)
        
        full_text += "\n\n--- End of Document ---\n\n"
    return full_text